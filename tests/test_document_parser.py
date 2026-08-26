from pptx import Presentation
from app.services.document_parser import parse_pptx


def test_pptx_text_extraction(tmp_path):
    path = tmp_path / "demo.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "AI in Education"
    slide.placeholders[1].text = "Personalized learning"
    prs.save(path)
    slides = parse_pptx(str(path))
    assert len(slides) == 1
    assert "AI in Education" in slides[0]["text"]
    assert "Personalized learning" in slides[0]["text"]
