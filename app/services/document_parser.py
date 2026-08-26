from pathlib import Path


def parse_pptx(file_path: str) -> list[dict]:
    from pptx import Presentation
    prs = Presentation(file_path)
    slides = []
    for number, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text:
                    texts.append(text)
        slides.append({"slide_number": number, "text": "\n".join(texts)})
    return slides


def parse_pdf(file_path: str) -> list[dict]:
    import pdfplumber
    pages = []
    with pdfplumber.open(file_path) as pdf:
        for number, page in enumerate(pdf.pages, start=1):
            pages.append({"slide_number": number, "text": (page.extract_text() or "").strip()})
    return pages


def parse_document(file_path: str, filename: str | None = None) -> list[dict]:
    name = (filename or Path(file_path).name).lower()
    if name.endswith(".pptx"):
        return parse_pptx(file_path)
    if name.endswith(".pdf"):
        return parse_pdf(file_path)
    raise ValueError("Unsupported presentation file. Use .pptx or .pdf.")
